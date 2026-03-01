import re
import io
import tiktoken
import pdfplumber

from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from typing import List, Tuple
from odf import text as odf_text
from sqlalchemy.orm import Session
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load as odf_load

from app.models.models import MaterialChunk, User
from app.core.exceptions import ValidationException
from app.repositories.chunk_repository import ChunkRepository

class ChunkService:
    def __init__(self, db: Session):
        self.chunk_repo = ChunkRepository(db)

    def chunk_text_semantic(self,
            self_text: str,
            model: str = "gpt-4o-mini",
            max_tokens: int = 600,
            overlap_tokens: int = 100,
        ) -> Tuple[List, List]:
        enc = tiktoken.encoding_for_model(model)

        sentences = re.split(r'(?<=[.!?])\s+', self_text)

        chunks = []
        current_chunk = []
        current_tokens = 0
        tokens = []

        for sentence in sentences:
            sentence_tokens = len(enc.encode(sentence))

            # if adding this sentence exceeds limit → flush
            if current_tokens + sentence_tokens > max_tokens and current_chunk:
                chunk_text = " ".join(current_chunk)
                chunks.append(chunk_text)
                tokens.append(current_tokens)

                # create overlap
                overlap_text = enc.encode(chunk_text)[-overlap_tokens:]
                overlap_str = enc.decode(overlap_text)

                current_chunk = [overlap_str, sentence]
                current_tokens = len(enc.encode(overlap_str)) + sentence_tokens
            else:
                current_chunk.append(sentence)
                current_tokens += sentence_tokens

        if current_chunk:
            chunks.append(" ".join(current_chunk))
            tokens.append(current_tokens)

        return chunks, tokens

    async def extract_and_chunk(self, file_data: dict, **kwargs) -> Tuple[List, List]:
        raw_text = await self.extract_text(**file_data)
        return self.chunk_text_semantic(raw_text, **kwargs)

    def create_chunks_for_material(self, user: User, material_id, chunks: Tuple[List, List]):
        material_chunks = []
        index = 0
        for chunk, token in zip(chunks[0], chunks[1]):
            material_chunk = MaterialChunk(
                user_id=user.user_id,
                material_id=material_id, 
                chunk_index=index,
                content=chunk,
                token_count=token
            )
            material_chunk = self.chunk_repo.create(material_chunk)
            material_chunks.append(material_chunk)
            index += 1
        return material_chunks

    async def extract_text(self, file_bytes: bytes,
                        content_type: str | None = None,
                        file_name: str | None = None) -> str:
        
        content_type = content_type or ""
        filename = (file_name or "").lower()

        # ---------- TEXT / MARKDOWN ----------
        if content_type.startswith("text/") or filename.endswith((".txt", ".md")):
            return self.clean_text(file_bytes.decode("utf-8", errors="ignore"))

        # ---------- PDF ----------
        if content_type == "application/pdf" or filename.endswith(".pdf"):
            return self.clean_text(self.extract_pdf_text(file_bytes))

        # ---------- DOCX ----------
        if (
            content_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            or filename.endswith(".docx")
        ):
            return self.clean_text(self.extract_docx_text(file_bytes))

        # ---------- DOC (legacy) ----------
        if content_type == "application/msword" or filename.endswith(".doc"):
            # basic fallback — many .doc files won't parse well
            return self.clean_text(file_bytes.decode("utf-8", errors="ignore"))

        # ---------- RTF ----------
        if content_type == "application/rtf" or filename.endswith(".rtf"):
            return self.clean_text(rtf_to_text(file_bytes.decode("utf-8", errors="ignore")))

        # ---------- PPT / PPTX ----------
        if (
            content_type
            in [
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ]
            or filename.endswith((".ppt", ".pptx"))
        ):
            return self.clean_text(self.extract_pptx_text(file_bytes))

        # ---------- ODT ----------
        if (
            content_type == "application/vnd.oasis.opendocument.text"
            or filename.endswith(".odt")
        ):
            return self.clean_text(self.extract_odt_text(file_bytes))

        # ---------- ODP ----------
        if (
            content_type == "application/vnd.oasis.opendocument.presentation"
            or filename.endswith(".odp")
        ):
            return self.clean_text(self.extract_odp_text(file_bytes))

        # ---------- HTML fallback ----------
        if content_type == "text/html" or filename.endswith(".html"):
            soup = BeautifulSoup(file_bytes, "html.parser")
            return self.clean_text(soup.get_text(separator="\n"))

        raise ValidationException("Unsupported file type for text extraction")

    def extract_pdf_text(self, file_bytes: bytes) -> str:
        texts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "\n".join(texts)
    
    def extract_docx_text(self, file_bytes: bytes) -> str:
        doc = Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text)

    def extract_pptx_text(self, file_bytes: bytes) -> str:
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text:
                        texts.append(shape.text)

        return "\n".join(texts)
    
    def extract_odp_text(self, file_bytes: bytes) -> str:
        doc = odf_load(io.BytesIO(file_bytes))
        texts = []

        for p in doc.getElementsByType(odf_text.P):
            texts.append(str(p))

        return "\n".join(texts)
        
    def extract_odt_text(self, file_bytes: bytes) -> str:
        doc = odf_load(io.BytesIO(file_bytes))
        texts = []

        for p in doc.getElementsByType(odf_text.P):
            texts.append(str(p))

        return "\n".join(texts)
    

    def clean_text(self, text: str) -> str:
        text = text.replace("\x00", " ")
        text = re.sub(r"\s+", " ", text)
        return text.strip()

    def get_chunks_by_material_id(self, material_id: str, **kwargs):
        return self.chunk_repo.get_by_material_id(material_id, **kwargs)

    # def delete_chunks_for_material(self, material_id):
    #     chunks = self.chunk_repo.get_by_material_id(material_id)
    #     for chunk in chunks:
    #         self.chunk_repo.delete(chunk.chunk_id)